from pptx import Presentation

prs = Presentation(r'C:\Users\paron\Downloads\Project Guide.pptx')

for i in [18]:
    slide = prs.slides[i]
    print('\n========== SLIDE %d ==========' % (i+1))
    for shape in slide.shapes:
        print('\nShape: %s, name=%s, pos=(%s, %s), size=(%s, %s)' % (shape.shape_type, shape.name, shape.left, shape.top, shape.width, shape.height))
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                runs = [r.text for r in para.runs]
                txt = ''.join(runs)
                if txt:
                    print('  Text: %s' % txt)
        if shape.has_table:
            table = shape.table
            print('  TABLE %dx%d:' % (len(table.rows), len(table.columns)))
            for r_idx, row in enumerate(table.rows):
                cells = [cell.text for cell in row.cells]
                print('    Row %d: %s' % (r_idx, cells))
